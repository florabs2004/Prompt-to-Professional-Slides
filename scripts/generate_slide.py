import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_ge_slide(data):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(19, 19, 20)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('headline', 'Strategic Overview')
    p.font.name = 'Google Sans'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Points (Bento Boxes)
    points = data.get('points', [])
    num_points = len(points)
    
    if num_points > 0:
        card_width = (13.333 - 1.0 - (0.5 * (num_points - 1))) / num_points
        for i, point in enumerate(points):
            left = 0.5 + (i * (card_width + 0.5))
            top = 2.5
            height = 4.0
            
            # Card Shape
            shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(card_width), Inches(height))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(30, 31, 32)
            shape.line.color.rgb = RGBColor(60, 60, 60)

            # Stat
            stat_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(card_width - 0.4), Inches(0.8))
            tf = stat_box.text_frame
            p = tf.paragraphs[0]
            p.text = point.get('stat', '')
            p.font.name = 'Google Sans'
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(138, 180, 248)

            # Title
            t_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 1.2), Inches(card_width - 0.4), Inches(0.5))
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = point.get('title', '')
            p.font.name = 'Google Sans'
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

            # Body
            b_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 1.8), Inches(card_width - 0.4), Inches(1.5))
            tf = b_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = point.get('body', '')
            p.font.name = 'Google Sans'
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(154, 160, 166)

    # Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Gemini Enterprise | Strategy Proposal"
    p.font.name = 'Google Sans'
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(154, 160, 166)

    prs.save(data.get('filename', 'strategy_slide.pptx'))

if __name__ == "__main__":
    content = json.loads(sys.argv[1])
    create_ge_slide(content)
