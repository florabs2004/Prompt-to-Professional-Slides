import sys
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def load_data():
    base_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    schema_path = os.path.join(base_path, 'assets', 'mapping_schema.json')
    content_path = os.path.join(base_path, 'assets', 'master_content.json')
    
    with open(schema_path, 'r') as f:
        schema = json.load(f)
    with open(content_path, 'r') as f:
        content = json.load(f)
    return schema, content

def add_slide_with_card(prs, headline, stat, title, body, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(19, 19, 20)

    # Slide Meta
    meta = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(0.5))
    tf = meta.text_frame
    p = tf.paragraphs[0]
    p.text = f"SECTION {slide_num} / GEMINI ENTERPRISE"
    p.font.name = 'Google Sans'
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(154, 160, 166)

    # Headline
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = headline
    p.font.name = 'Google Sans'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Card
    left, top, width, height = Inches(0.5), Inches(2.5), Inches(6.0), Inches(4.0)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 31, 32)
    shape.line.color.rgb = RGBColor(60, 60, 60)

    # Stat
    stat_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), width - Inches(0.6), Inches(1.0))
    p = stat_box.text_frame.paragraphs[0]
    p.text = stat
    p.font.name = 'Google Sans'
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(138, 180, 248)

    # Card Title
    t_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(1.5), width - Inches(0.6), Inches(0.5))
    p = t_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = 'Google Sans'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Card Body
    b_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(2.2), width - Inches(0.6), Inches(1.5))
    tf = b_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = body
    p.font.name = 'Google Sans'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(154, 160, 166)

def generate_deck(params):
    schema, content = load_data()
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # 1. Intro Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(19, 19, 20)
    
    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12), Inches(2.0))
    p = t_box.text_frame.paragraphs[0]
    p.text = f"Google Cloud:\n{params.get('industry', 'Enterprise')} Transformation"
    p.font.name = 'Google Sans'
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 2. Industry Context
    ind_id = schema['industries'].get(params.get('industry'), 'ind_horizontal')
    ind_data = content['industry_context'].get(ind_id)
    add_slide_with_card(prs, ind_data['headline'], ind_data['stat'], ind_data['title'], ind_data['body'], 1)

    # 3. Tech Deep-Dive
    tech_id = schema['tech'].get(params.get('tech'), 'tech_horizontal')
    tech_data = content['tech_deep_dive'].get(tech_id)
    add_slide_with_card(prs, tech_data['headline'], tech_data['stat'], tech_data['title'], tech_data['body'], 2)

    # 4. LoB Value
    lob_id = schema['lobs'].get(params.get('lob'), 'lob_horizontal')
    lob_data = content['lob_value'].get(lob_id)
    add_slide_with_card(prs, lob_data['headline'], lob_data['stat'], lob_data['title'], lob_data['body'], 3)

    # 5. Next Steps
    inc_id = schema['incentives'].get(params.get('incentive'), 'inc_standard')
    inc_data = content['next_steps'].get(inc_id)
    add_slide_with_card(prs, inc_data['headline'], inc_data['stat'], inc_data['title'], inc_data['body'], 4)

    prs.save('customized_pitch_deck.pptx')

if __name__ == "__main__":
    params = json.loads(sys.argv[1])
    generate_deck(params)
